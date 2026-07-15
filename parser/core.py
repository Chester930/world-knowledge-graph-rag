import io
import logging
import os
import re
import tempfile
from pathlib import Path
from typing import List, Optional

try:
    from parser.image_pipeline import (
        ImagePipeline, ImagePipelineConfig, _CAPTION_PATTERNS, emu_to_px, points_to_px,
    )
except ImportError:
    from image_pipeline import (
        ImagePipeline, ImagePipelineConfig, _CAPTION_PATTERNS, emu_to_px, points_to_px,
    )

logger = logging.getLogger(__name__)

# 第三方庫導入與容錯處理
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from pdfminer.high_level import extract_pages
    from pdfminer.layout import LTTextContainer, LTChar, LTFigure, LTImage
except ImportError:
    extract_pages = None
    LTFigure = None
    LTImage = None

try:
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn as docx_qn
except ImportError:
    DocxDocument = None
    docx_qn = None

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None
    MSO_SHAPE_TYPE = None

try:
    from pdf2image import convert_from_path
    import pytesseract
    from PIL import Image
    # Windows 預設 tesseract 安裝路徑檢查，以防 PATH 未設定
    TESSERACT_CMD_CANDIDATES = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        r"C:\Users\666\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"
    ]
    for cand in TESSERACT_CMD_CANDIDATES:
        if os.path.exists(cand):
            pytesseract.pytesseract.tesseract_cmd = cand
            break
except ImportError:
    convert_from_path = None
    pytesseract = None

try:
    import trafilatura
except ImportError:
    trafilatura = None

try:
    from youtube_transcript_api import YouTubeTranscriptApi
except ImportError:
    YouTubeTranscriptApi = None

try:
    import whisper
except ImportError:
    whisper = None

try:
    import yt_dlp
except ImportError:
    yt_dlp = None


_whisper_model_cache = {}


def _transcribe_with_whisper(path: str, model_size: str = "tiny") -> str:
    """共用的本地 Whisper 語音轉文字函式，供音檔上傳與 YouTube 音軌備援共用。

    模型載入後會快取在行程記憶體中，避免重複解析時重複載入。
    """
    if whisper is None:
        raise ImportError("未安裝 openai-whisper 套件，請執行 pip install openai-whisper")

    if model_size not in _whisper_model_cache:
        _whisper_model_cache[model_size] = whisper.load_model(model_size)
    model = _whisper_model_cache[model_size]

    result = model.transcribe(str(path))
    return result.get("text", "").strip()


def _iter_lt_images(item):
    """遞迴尋找 LTFigure 容器內的所有 LTImage 元素（LTFigure 可能巢狀包含子圖形）。"""
    if LTImage is not None and isinstance(item, LTImage):
        yield item
        return
    try:
        children = list(item)
    except TypeError:
        return
    for child in children:
        yield from _iter_lt_images(child)


def _pil_from_ltimage(lt_image):
    """將 pdfminer 的 LTImage 物件盡力解碼為 PIL Image。

    多數內嵌影像為 JPEG（DCTDecode），pdfminer 的 get_data() 對此類濾鏡不會二次解碼，
    可直接交給 PIL 開啟；其餘濾鏡（如未壓縮點陣圖）則嘗試以 srcsize/bits 還原像素資料。
    無法解碼時回傳 None，呼叫端會安全略過該圖片，不影響其餘解析流程。
    """
    if Image is None:
        return None
    if getattr(lt_image, "imagemask", False):
        return None  # 遮罩用途的裝飾性影像（如底色色塊），非內容圖片

    try:
        data = lt_image.stream.get_data()
    except Exception:
        return None

    try:
        img = Image.open(io.BytesIO(data))
        img.load()
        return img.convert("RGB")
    except Exception:
        pass

    try:
        width, height = lt_image.srcsize
        bits = getattr(lt_image, "bits", 8)
        if bits == 8 and len(data) >= width * height * 3:
            return Image.frombytes("RGB", (width, height), data[: width * height * 3])
        if bits == 8 and len(data) >= width * height:
            return Image.frombytes("L", (width, height), data[: width * height]).convert("RGB")
    except Exception:
        pass

    return None


class DocumentParserError(Exception):
    """文件解析異常基類"""
    pass


class DocumentParser:
    """統一文件解析進入點"""

    def __init__(self, image_config: Optional[ImagePipelineConfig] = None):
        """image_config 為圖文統一處理管線的設定；不提供則使用預設值
        （OCR 核心能力預設開啟、圖理解模型預設關閉），維持向後相容。"""
        self.image_config = image_config or ImagePipelineConfig()
        self.image_pipeline = ImagePipeline(self.image_config)

    def parse_file(self, file_path: str) -> str:
        """根據檔案副檔名進行路由解析，返回純文字內容。"""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"找不到檔案: {file_path}")

        self.image_pipeline.reset_numbering()

        suffix = path.suffix.lower()
        if suffix in [".txt", ".md"]:
            return self._parse_text(path)
        elif suffix == ".pdf":
            return self._parse_pdf(path)
        elif suffix in [".docx", ".doc"]:
            return self._parse_docx(path)
        elif suffix in [".pptx", ".ppt"]:
            return self._parse_pptx(path)
        elif suffix in [".mp3", ".wav", ".m4a", ".ogg", ".flac", ".mp4", ".mkv", ".webm"]:
            return self._parse_audio(path)
        else:
            raise DocumentParserError(f"不支援的檔案格式: {suffix}")

    def _parse_text(self, path: Path) -> str:
        """解析 TXT 或 Markdown"""
        try:
            # 優先嘗試 UTF-8，失敗則嘗試 Big5 (繁體中文常見編碼) 或 cp950
            return path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            try:
                return path.read_text(encoding="big5")
            except UnicodeDecodeError:
                return path.read_text(encoding="cp950", errors="ignore")

    def _parse_docx(self, path: Path) -> str:
        """解析 Word 文件，保留表格 Markdown 格式，並依段落順序抽取內嵌圖片"""
        if DocxDocument is None:
            raise ImportError("未安裝 python-docx 套件，請執行 pip install python-docx")

        doc = DocxDocument(path)
        content_parts = []

        # DOCX 為流式版面、無座標資訊，改以文件既有段落順序作為「閱讀順序」的替代依據
        full_doc_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        # 純文字預檢：整份文件的關聯清單只需掃描一次，即可判斷是否存在任何內嵌圖片；
        # 純文字文件（無圖片）可完全跳過逐段落的 XML 掃描，避免對每個段落重複做無意義的查找。
        has_images = self.image_config.enable_ocr and self._docx_has_embedded_images(doc)

        # 遍歷 docx 的所有子元素，以保留表格與段落的原始順序
        for element in doc.element.body:
            # 如果是段落
            if element.tag.endswith('p'):
                # 尋找對應的 python-docx Paragraph 對象
                for p in doc.paragraphs:
                    if p._element is element:
                        if p.text.strip():
                            content_parts.append(p.text)
                        if has_images:
                            content_parts.extend(
                                self._process_docx_paragraph_images(doc, p, full_doc_text)
                            )
                        break
            # 如果是表格
            elif element.tag.endswith('tbl'):
                for t in doc.tables:
                    if t._element is element:
                        # 轉為 Markdown table
                        table_md = self._convert_table_to_markdown(t)
                        if table_md:
                            content_parts.append("\n" + table_md + "\n")
                        break

        return "\n\n".join(content_parts)

    def _docx_has_embedded_images(self, doc) -> bool:
        """輕量預檢：掃描一次文件關聯清單（rels）判斷是否存在任何內嵌圖片，
        取代逐段落做 XML findall 的重複查找。任何解析失敗都保守回傳 True，
        交由後續逐段落掃描處理，避免誤判漏圖。"""
        try:
            return any("image" in rel.reltype.lower() for rel in doc.part.rels.values())
        except Exception:
            return True

    def _process_docx_paragraph_images(self, doc, paragraph, full_doc_text: str) -> List[str]:
        """抽取段落中內嵌的圖片（DrawingML a:blip），逐一經圖片管線處理，
        回傳格式化文字區塊清單；任何失敗都安全略過，不中斷整份文件解析。"""
        if Image is None or docx_qn is None:
            return []

        try:
            blips = paragraph._element.findall('.//' + docx_qn('a:blip'))
        except Exception:
            return []

        if not blips:
            return []

        caption_text = paragraph.text.strip()
        if not caption_text:
            # Word「插入標題」慣例：圖片自成一段（通常無文字），標題文字在緊接的下一段。
            # 優先向下找 1 個非空段落，其次向上找 1 個非空段落，僅在該段文字符合圖號
            # 標註格式（如「圖1 系統架構」）時才採用，避免誤把無關段落當成圖說。
            try:
                p_idx = doc.paragraphs.index(paragraph)
                for offset in (1, -1):
                    target_idx = p_idx + offset
                    if 0 <= target_idx < len(doc.paragraphs):
                        candidate_text = doc.paragraphs[target_idx].text.strip()
                        if candidate_text and any(pat.match(candidate_text) for pat in _CAPTION_PATTERNS):
                            caption_text = candidate_text
                            break
            except Exception:
                pass

        blocks = []
        for blip in blips:
            rId = blip.get(docx_qn('r:embed'))
            if not rId:
                continue
            try:
                image_part = doc.part.related_parts[rId]
                pil_image = Image.open(io.BytesIO(image_part.blob)).convert("RGB")
            except Exception:
                continue

            result = self.image_pipeline.process_image(
                pil_image,
                doc_type="docx",
                nearby_caption_text=caption_text,
                full_doc_text=full_doc_text,
                display_size_px=self._docx_image_display_size_px(blip),
            )
            block = result.to_text_block() if result else ""
            if block:
                blocks.append(block)

        return blocks

    def _docx_image_display_size_px(self, blip):
        """從 `a:blip` 元素往上尋找其所屬的 `wp:inline`/`wp:anchor`，取得裡面的
        `wp:extent`（該圖片實際顯示尺寸，EMU 單位，與圖片檔案本身的原始像素解析度無關）
        並換算成 96 DPI 基準像素。DOCX 雖然沒有頁面座標，但插入圖片時一定會指定顯示尺寸，
        此資訊即來自這裡。找不到（結構異常）時回傳 None，由 process_image 安全退回使用
        原始像素尺寸判斷，不影響正確性。"""
        if docx_qn is None:
            return None
        try:
            node = blip.getparent()
            while node is not None:
                if node.tag in (docx_qn('wp:inline'), docx_qn('wp:anchor')):
                    extent = node.find(docx_qn('wp:extent'))
                    if extent is not None:
                        cx, cy = extent.get('cx'), extent.get('cy')
                        if cx and cy:
                            return (emu_to_px(float(cx)), emu_to_px(float(cy)))
                    break
                node = node.getparent()
        except Exception:
            pass
        return None

    def _convert_table_to_markdown(self, table) -> str:
        """將 python-docx 的 Table 物件轉換為 Markdown 表格語法，保留儲存格內換行"""
        if not table.rows:
            return ""

        md_rows = []
        # 處理標頭列
        header_cells = table.rows[0].cells
        md_rows.append("| " + " | ".join(cell.text.replace("\n", "<br>").strip() for cell in header_cells) + " |")
        md_rows.append("| " + " | ".join("---" for _ in header_cells) + " |")

        # 處理資料列
        for row in table.rows[1:]:
            cells = row.cells
            row_texts = []
            prev_cell = None
            for cell in cells:
                if cell == prev_cell:
                    # 跨欄合併，留空
                    row_texts.append("")
                else:
                    # 將內建換行替換為 <br> 以維持 Markdown 表格單行結構
                    row_texts.append(cell.text.replace("\n", "<br>").strip())
                prev_cell = cell
            md_rows.append("| " + " | ".join(row_texts) + " |")

        return "\n".join(md_rows)

    def _parse_pptx(self, path: Path) -> str:
        """解析 PPTX 投影片，並依規格抽取內嵌圖片（PPTX 屬視覺導向文件類型，
        整頁大圖/流程圖機率高，圖片管線的文件類型基礎分較高）。"""
        if Presentation is None:
            raise ImportError("未安裝 python-pptx 套件，請執行 pip install python-pptx")

        prs = Presentation(path)
        # 純文字預檢：逐張投影片掃描其自身的圖片關聯（PPTX 圖片關聯存在於各 slide
        # 自己的 rels，並非簡報層級的 prs.part.rels），純文字投影片可完全跳過逐 shape
        # 的圖片判斷與處理。
        has_images = self.image_config.enable_ocr and self._pptx_has_embedded_images(prs)
        content_parts = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"--- [Slide {slide_idx + 1}] ---")

            # 先收集本張投影片所有文字，供圖片的「正文明確圖號引用」比對使用
            full_slide_text = "\n".join(
                shape.text.strip() for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )

            # 依 (top, left) 座標重建視覺閱讀順序，而非沿用 slide.shapes 的
            # z-order/建立順序，確保圖片鄰近文字（last_text_seen）與最終輸出順序
            # 符合投影片實際版面配置。
            ordered_shapes = self._pptx_reading_order(list(slide.shapes))

            last_text_seen = ""
            for shape in ordered_shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    slide_text.append(text)
                    last_text_seen = text
                elif (
                    has_images
                    and MSO_SHAPE_TYPE is not None
                    and shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                ):
                    image_block = self._process_pptx_picture_shape(shape, last_text_seen, full_slide_text)
                    if image_block:
                        slide_text.append(image_block)

            content_parts.append("\n".join(slide_text))

        return "\n\n".join(content_parts)

    def _pptx_has_embedded_images(self, prs) -> bool:
        """輕量預檢：逐張投影片掃描其自身的關聯清單（slide.part.rels）判斷是否存在任何
        內嵌圖片。PPTX 的圖片關聯存在於各 slide 自己的 rels，並非簡報層級的
        `prs.part.rels`（後者僅含投影片母片/版面配置/主題等簡報層級關聯，不含各投影片
        內的圖片）。任何解析失敗都保守回傳 True，交由逐 shape 判斷處理，避免誤判漏圖。"""
        if prs is None:
            return True
        try:
            for slide in prs.slides:
                try:
                    for rel in slide.part.rels.values():
                        if "image" in rel.reltype.lower():
                            return True
                except Exception:
                    continue
            return False
        except Exception:
            return True

    def _pptx_reading_order(self, shapes: list) -> list:
        """依 (top, left) EMU 座標重建投影片 shape 的視覺閱讀順序：先依 top 座標
        排序後分列（容忍度內視為同一列），列內再依 left 由左到右排序。

        採兩階段分桶排序而非成對比較器排序，避免「跨列容忍度」判斷在鏈式比較時
        不滿足遞移律，導致排序結果不穩定的問題。
        """
        ROW_TOLERANCE_EMU = 228600  # 0.25 inch，同一橫列的垂直座標容忍範圍

        items = [(getattr(s, "top", 0) or 0, getattr(s, "left", 0) or 0, s) for s in shapes]
        items.sort(key=lambda it: it[0])

        rows = []
        current_row = []
        current_row_top = None
        for top, left, shape in items:
            if current_row and abs(top - current_row_top) > ROW_TOLERANCE_EMU:
                rows.append(current_row)
                current_row = []
                current_row_top = None
            current_row.append((top, left, shape))
            if current_row_top is None:
                current_row_top = top
        if current_row:
            rows.append(current_row)

        ordered = []
        for row in rows:
            row.sort(key=lambda it: it[1])
            ordered.extend(shape for _, _, shape in row)
        return ordered

    def _process_pptx_picture_shape(self, shape, nearby_text: str, full_slide_text: str) -> str:
        """處理單一 PPTX 圖片 shape，回傳格式化後的文字區塊；任何失敗都安全降級為空字串略過。"""
        if Image is None:
            return ""
        try:
            pil_image = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
        except Exception:
            return ""

        display_size_px = None
        try:
            display_size_px = (emu_to_px(shape.width), emu_to_px(shape.height))
        except Exception:
            pass  # 取不到顯示尺寸時，process_image 會安全退回原始像素尺寸判斷

        result = self.image_pipeline.process_image(
            pil_image,
            doc_type="pptx",
            nearby_caption_text=nearby_text,
            full_doc_text=full_slide_text,
            display_size_px=display_size_px,
        )
        return result.to_text_block() if result else ""

    def _parse_audio(self, path: Path) -> str:
        """使用本地 Whisper 模型進行語音轉文字"""
        try:
            # 使用 tiny 模型，以防使用者下載過久。預設對中文已堪用
            return _transcribe_with_whisper(str(path), model_size="tiny")
        except ImportError:
            raise
        except Exception as e:
            raise DocumentParserError(f"音訊轉譯失敗: {str(e)}")

    def _parse_pdf(self, path: Path) -> str:
        """三層備援 PDF 轉譯主入口"""
        # 開啟一次 PdfReader，供軌道一與純文字預檢共用，避免重複開檔。
        # 開啟失敗（含未安裝 pypdf）都優雅降級為 None，讓後續軌道二/三照常運作，
        # 不因缺少 pypdf 而讓整個 PDF 解析直接失敗。
        reader = None
        if PdfReader is not None:
            try:
                reader = PdfReader(path)
            except Exception:
                reader = None

        # --- 軌道一：pypdf 快速文本流提取 ---
        text = self._pdf_track_pypdf(reader)
        used_ocr_fallback = False
        pages_layout = None

        # 評估是否需要 Fallback：字元數過少或疑似亂碼
        if self._is_low_quality_text(text):
            # --- 軌道二：pdfminer 佈局感知與雙欄排序提取 ---
            # 版面樹狀結構載入後會保留供圖文管線共用，避免同一份 PDF 被 extract_pages() 解析兩次。
            pages_layout = self._load_pdf_layout(path)
            text = self._pdf_track_pdfminer(pages_layout)

            # 如果依然判定為低品質（例如圖片掃描件），進入軌道三
            if self._is_low_quality_text(text):
                # --- 軌道三：OCR 視覺區域識別解析 ---
                text = self._pdf_track_ocr(path)
                used_ocr_fallback = True

        # --- 圖文統一處理：內嵌圖片抽取（核心能力，獨立於前述文字提取軌道結果） ---
        # 全頁已透過軌道三整頁 OCR 的情況下，內嵌圖片內容通常已被涵蓋於整頁 OCR 結果中，
        # 故僅在軌道一/二成功時才另行處理內嵌圖片，避免重複解析同一份掃描影像。
        #
        # 純文字預檢：先輕量檢查頁面資源中是否存在任何內嵌圖片 XObject，這比 pdfminer 的
        # extract_pages() 完整版面重建便宜得多。純文字文件（如論文、合約）通常完全沒有
        # 內嵌圖片，可直接跳過整套圖片管線，避免白白重新解析一次全文件版面。
        if self.image_config.enable_ocr and not used_ocr_fallback and self._pdf_has_embedded_images(reader if reader is not None else path):
            if pages_layout is None:
                pages_layout = self._load_pdf_layout(path)
            image_text = self._extract_pdf_images_text(pages_layout)
            if image_text:
                text = text + "\n\n" + image_text

        return text

    def _pdf_has_embedded_images(self, path_or_reader) -> bool:
        """輕量預檢：僅讀取各頁 XObject 資源清單判斷是否含任何內嵌圖片，
        不做 pdfminer 的座標/版面重建、也不做 pypdf 較重的影像解碼，成本遠低於
        `_extract_pdf_images_text`。可傳入檔案路徑（會自行開啟 PdfReader）或已開啟的
        PdfReader 實例（供 `_parse_pdf` 內部共用軌道一已開啟的 reader，避免重複開檔）。
        任何解析失敗都保守回傳 True，交由後續完整管線處理，避免誤判漏圖。"""
        if PdfReader is None:
            return True
        try:
            reader = path_or_reader if isinstance(path_or_reader, PdfReader) else PdfReader(path_or_reader)
        except Exception:
            return True

        try:
            for page in reader.pages:
                try:
                    resources = page.get("/Resources")
                    if not resources:
                        continue
                    xobjects = resources.get("/XObject")
                    if not xobjects:
                        continue
                    for obj_name in xobjects:
                        try:
                            if xobjects[obj_name].get("/Subtype") == "/Image":
                                return True
                        except Exception:
                            continue
                except Exception:
                    continue
            return False
        except Exception:
            return True

    def _pdf_track_pypdf(self, reader) -> str:
        """第一軌：使用已開啟的 PdfReader 提取文字"""
        if reader is None:
            return ""
        try:
            pages_text = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    pages_text.append(t)
            return "\n\n".join(pages_text)
        except Exception:
            return ""

    def _load_pdf_layout(self, path: Path) -> list:
        """載入 pdfminer 版面樹狀結構（LTPage 列表），供軌道二文字重建與圖片管線共用，
        避免同一份 PDF 被 `extract_pages()` 重複解析兩次。"""
        if extract_pages is None:
            return []
        try:
            return list(extract_pages(path))
        except Exception:
            return []

    def _pdf_track_pdfminer(self, pages_layout: list) -> str:
        """第二軌：使用已載入的版面樹重建閱讀順序與雙欄解析 (混合佈局優化版)"""
        if not pages_layout:
            return ""
        try:
            pages_text = []
            for page_layout in pages_layout:
                page_width = page_layout.width
                mid_x = page_width / 2

                # 提取該頁所有文字容器
                text_containers = []
                for element in page_layout:
                    if isinstance(element, LTTextContainer):
                        text_containers.append(element)

                if not text_containers:
                    continue

                # 1. 混合分流：將「跨欄大標題/全寬物件」與「窄欄正文」區分開來
                span_elements = []  # 跨中線的大型元素 (如 Title)
                column_elements = []  # 需要雙欄分流排序的元素

                for elem in text_containers:
                    width = elem.x1 - elem.x0
                    # 如果寬度大於頁面寬度的 70%，或者明顯橫跨中線且寬度寬
                    if width > (page_width * 0.65) and (elem.x0 < mid_x < elem.x1):
                        span_elements.append(elem)
                    else:
                        column_elements.append(elem)

                # 2. 對需要雙欄分流的元素進行左右分類
                left_col = []
                right_col = []
                for elem in column_elements:
                    elem_mid_x = (elem.x0 + elem.x1) / 2
                    if elem_mid_x < mid_x:
                        left_col.append(elem)
                    else:
                        right_col.append(elem)

                # 3. 各自按 Y1 軸降序（自上而下）排序
                span_elements.sort(key=lambda e: e.y1, reverse=True)
                left_col.sort(key=lambda e: e.y1, reverse=True)
                right_col.sort(key=lambda e: e.y1, reverse=True)

                # 4. 拓撲結構重組：按 Y 軸高度，將單欄元素插回雙欄元素之間
                page_content = []
                left_idx = 0
                right_idx = 0

                # 依高度從上到下依次處理單欄跨欄元素
                for span_elem in span_elements:
                    # 先將高度高於此跨欄元素的所有左右欄元素寫入
                    while left_idx < len(left_col) and left_col[left_idx].y1 > span_elem.y1:
                        page_content.append(left_col[left_idx].get_text().strip())
                        left_idx += 1
                    while right_idx < len(right_col) and right_col[right_idx].y1 > span_elem.y1:
                        page_content.append(right_col[right_idx].get_text().strip())
                        right_idx += 1
                    
                    # 寫入跨欄元素本身 (如論文標題或大圖表)
                    page_content.append(span_elem.get_text().strip())

                # 將剩餘的低於最後一個跨欄元素的所有左右欄元素寫入
                while left_idx < len(left_col):
                    page_content.append(left_col[left_idx].get_text().strip())
                    left_idx += 1
                while right_idx < len(right_col):
                    page_content.append(right_col[right_idx].get_text().strip())
                    right_idx += 1

                pages_text.append("\n".join(page_content))

            return "\n\n".join(pages_text)
        except Exception:
            return ""

    def _pdf_track_ocr(self, path: Path) -> str:
        """第三軌：PDF 轉影像後使用 pytesseract 進行 OCR"""
        if convert_from_path is None or pytesseract is None:
            raise ImportError(
                "未安裝 pdf2image 或 pytesseract 套件，或系統缺乏 Poppler/Tesseract。"
                "請執行 pip install pdf2image pytesseract 並安裝對應系統工具。"
            )

        try:
            # 將 PDF 所有頁面轉換為 PIL Image 物件
            images = convert_from_path(path)
            ocr_texts = []
            
            for idx, img in enumerate(images):
                # 提取繁體中文 (chi_tra) 和英文 (eng)
                try:
                    text = pytesseract.image_to_string(img, lang="chi_tra+eng")
                except Exception:
                    # 若本機未安裝繁體中文語言包，回退到英文與通用簡體中文
                    text = pytesseract.image_to_string(img, lang="eng")
                
                if text.strip():
                    ocr_texts.append(f"--- [Page {idx + 1} OCR] ---\n" + text.strip())
            
            return "\n\n".join(ocr_texts)
        except Exception as e:
            raise DocumentParserError(f"PDF OCR 解析失敗: {str(e)}")

    def _extract_pdf_images_text(self, pages_layout: list) -> str:
        """圖文統一處理管線：從已載入的版面樹狀結構抽取 PDF 內嵌圖片，經空間去重、評分、
        編號、（選配）圖理解模型後，依頁面由上到下的閱讀順序整理為文字區塊。

        任何非預期錯誤都會被吞下並記錄警告，絕不讓圖片處理中斷主要的文字解析流程。
        """
        if not pages_layout or Image is None:
            return ""

        page_blocks = []
        try:
            for page_idx, page_layout in enumerate(pages_layout):
                text_containers = []
                figures = []
                for element in page_layout:
                    if isinstance(element, LTTextContainer):
                        text_containers.append(element)
                    elif LTFigure is not None and isinstance(element, LTFigure):
                        figures.append(element)

                if not figures:
                    continue

                text_bboxes = [(c.x0, c.y0, c.x1, c.y1) for c in text_containers]
                page_full_text = "\n".join(c.get_text() for c in text_containers)
                doc_type = "pdf_native" if text_containers else "pdf_scanned"

                # 第一階段：收集所有存活圖片（通過空間去重與解碼），先依 y1 由上到下排序，
                # 才進入圖片管線指派圖號——確保圖號指派順序與視覺閱讀順序一致，
                # 而非依 pdfminer 內部圖形元素的（近似繪製順序的）發現順序指派。
                candidates = []
                for fig in figures:
                    for lt_image in _iter_lt_images(fig):
                        bbox = (lt_image.x0, lt_image.y0, lt_image.x1, lt_image.y1)

                        # 空間去重：原生文字已大面積覆蓋此區域，屬零成本可跳過
                        if self.image_pipeline.is_covered_by_native_text(bbox, text_bboxes):
                            continue

                        pil_image = _pil_from_ltimage(lt_image)
                        if pil_image is None:
                            continue

                        candidates.append((bbox, pil_image))

                if not candidates:
                    continue

                candidates.sort(key=lambda item: item[0][3], reverse=True)  # y1 由高到低=由上到下

                # 第二階段：依排序後順序逐一交給圖片管線處理（含圖號指派），
                # 輸出區塊順序與圖號指派順序自然一致。
                blocks = []
                for bbox, pil_image in candidates:
                    nearby_caption = self._find_nearby_pdf_caption(bbox, text_containers)
                    display_size_px = (
                        points_to_px(bbox[2] - bbox[0]),
                        points_to_px(bbox[3] - bbox[1]),
                    )
                    result = self.image_pipeline.process_image(
                        pil_image,
                        doc_type=doc_type,
                        nearby_caption_text=nearby_caption,
                        full_doc_text=page_full_text,
                        display_size_px=display_size_px,
                    )
                    block = result.to_text_block() if result else ""
                    if block:
                        blocks.append(block)
                if blocks:
                    page_blocks.append(f"--- [第 {page_idx + 1} 頁圖片解析] ---\n" + "\n\n".join(blocks))

        except Exception as e:
            logger.warning("[DocumentParser] PDF 圖片處理發生非預期錯誤，已略過圖片解析: %s", str(e))
            return ""

        return "\n\n".join(page_blocks)

    def _find_nearby_pdf_caption(
        self, image_bbox, text_containers: List[LTTextContainer], max_gap: float = 40.0
    ) -> str:
        """在圖片正上方或正下方、水平方向有重疊的文字區塊中尋找最接近者，作為圖說候選文字。"""
        ix0, iy0, ix1, iy1 = image_bbox
        best_text = ""
        best_gap = max_gap

        for c in text_containers:
            cx0, cy0, cx1, cy1 = c.x0, c.y0, c.x1, c.y1
            if cx1 < ix0 or cx0 > ix1:
                continue  # 水平方向無重疊，非同一欄的圖說候選

            if cy1 <= iy0:
                gap = iy0 - cy1  # 文字在圖片下方
            elif cy0 >= iy1:
                gap = cy0 - iy1  # 文字在圖片上方
            else:
                continue  # 與圖片重疊，通常是圖片內部疊字，非圖說

            if gap < best_gap:
                best_gap = gap
                best_text = c.get_text().strip()

        return best_text

    def _detect_double_column(self, containers: List[LTTextContainer], mid_x: float) -> bool:
        """啟發式偵測頁面是否為雙欄排版"""
        if len(containers) < 4:
            return False

        left_side_count = 0
        right_side_count = 0
        span_mid_count = 0

        for elem in containers:
            # 寬度太窄的可能是頁碼或小標，不列入判定
            if elem.x1 - elem.x0 < 30:
                continue
            
            # 判斷是否跨越中線
            if elem.x0 < mid_x and elem.x1 > mid_x:
                # 如果中線穿過文字塊，且文字塊寬度大於中線的 20%，視為跨越
                overlap = min(elem.x1, mid_x) - max(elem.x0, mid_x)
                if (elem.x1 - elem.x0) > (mid_x * 0.4):
                    span_mid_count += 1
            elif elem.x1 <= mid_x:
                left_side_count += 1
            elif elem.x0 >= mid_x:
                right_side_count += 1

        # 如果跨越中線的區塊極少，而兩側各有顯著數量的區塊，判定為雙欄
        if span_mid_count <= 2 and left_side_count >= 2 and right_side_count >= 2:
            return True
        return False

    def _is_low_quality_text(self, text: str) -> bool:
        """判定文字是否低品質（如空白、純符號、亂碼或字元數極低）"""
        clean_text = text.strip()
        if not clean_text:
            return True
        
        # 1. 字元數過少（整份文件少於 80 字，通常是掃描件）
        if len(clean_text) < 80:
            return True
            
        # 2. 亂碼與特殊符號比例過高檢測
        # 計算正常可讀中英文字元比例
        alphanumeric_chinese = len(re.findall(r'[\u4e00-\u9fa5a-zA-Z0-9]', clean_text))
        ratio = alphanumeric_chinese / len(clean_text)
        
        # 如果可讀字元比例低於 60%（表示有大量特殊控制碼或亂碼），判定為低品質
        if ratio < 0.6:
            return True
            
        return False


def sentence_aware_chunking(text: str, chunk_size: int = 500, chunk_overlap: int = 50) -> List[str]:
    """句子感知的文本切片演算法，確保句子邊界不被粗暴截斷"""
    # 句尾結束符號規則，包含中英文標點
    sentence_endings = re.compile(r'([。！？…\n]|\.\s|\?\s|!\s)')
    
    # 依句尾標點進行文本分割，保留標點符號
    parts = []
    current_pos = 0
    for match in sentence_endings.finditer(text):
        end_pos = match.end()
        parts.append(text[current_pos:end_pos])
        current_pos = end_pos
    if current_pos < len(text):
        parts.append(text[current_pos:])

    chunks = []
    current_chunk = []
    current_len = 0

    for part in parts:
        part_len = len(part)
        # 如果單句長度就大於 chunk_size，直接獨立為一個 chunk
        if part_len >= chunk_size:
            if current_chunk:
                chunks.append("".join(current_chunk))
                current_chunk = []
                current_len = 0
            chunks.append(part)
            continue

        if current_len + part_len > chunk_size:
            # 建立當前 chunk
            chunks.append("".join(current_chunk))
            # 保留重疊區段：從 current_chunk 的尾端抓取約 chunk_overlap 長度的句子
            overlap_parts = []
            overlap_len = 0
            for p in reversed(current_chunk):
                if overlap_len + len(p) <= chunk_overlap:
                    overlap_parts.insert(0, p)
                    overlap_len += len(p)
                else:
                    break
            current_chunk = overlap_parts + [part]
            current_len = overlap_len + part_len
        else:
            current_chunk.append(part)
            current_len += part_len

    if current_chunk:
        chunks.append("".join(current_chunk))

    return [c.strip() for c in chunks if c.strip()]


class URLParser:
    """網頁與 YouTube 連結解析器"""

    def parse_url(self, url: str) -> str:
        """解析 URL，支援一般網頁與 YouTube 影片"""
        url = url.strip()
        if not url:
            raise ValueError("URL 不能為空")

        if self._is_youtube_url(url):
            return self._parse_youtube(url)
        else:
            return self._parse_webpage(url)

    def _is_youtube_url(self, url: str) -> bool:
        """判斷是否為 YouTube 網址"""
        pattern = r'(https?://)?(www\.)?(youtube\.com|youtu\.be)/'
        return bool(re.match(pattern, url, re.IGNORECASE))

    def _parse_youtube(self, url: str) -> str:
        """提取 YouTube 字幕；若無字幕則自動備援下載音軌並以 Whisper 轉譯"""
        video_id = self._extract_youtube_id(url)
        if not video_id:
            raise ValueError(f"無法解析 YouTube 影片 ID: {url}")

        subtitle_error = None
        if YouTubeTranscriptApi is not None:
            try:
                full_text = self._fetch_youtube_subtitles(video_id)
                return f"--- [YouTube Video {video_id} 字幕逐字稿] ---\n\n" + full_text
            except Exception as e:
                subtitle_error = e
        else:
            subtitle_error = ImportError("未安裝 youtube-transcript-api 套件")

        # 備援軌：該影片無字幕（或字幕抓取失敗），改為下載音軌並用本地 Whisper 轉譯
        if yt_dlp is None:
            raise DocumentParserError(
                f"YouTube 字幕提取失敗: {str(subtitle_error)}\n"
                "(且未安裝 yt-dlp 套件，無法備援下載音軌，請執行 pip install yt-dlp)"
            )

        try:
            audio_text = self._transcribe_youtube_audio(url, video_id)
            return f"--- [YouTube Video {video_id} 音軌 Whisper 逐字稿（無可用字幕，已自動備援轉譯）] ---\n\n" + audio_text
        except Exception as e:
            raise DocumentParserError(
                f"YouTube 字幕提取失敗: {str(subtitle_error)}\n音軌備援轉譯亦失敗: {str(e)}"
            )

    def _fetch_youtube_subtitles(self, video_id: str) -> str:
        """嘗試抓取 YouTube 官方或社群字幕，找不到則拋出例外交由呼叫端處理備援"""
        languages = ['zh-TW', 'zh-HK', 'zh-CN', 'zh', 'en']
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)

        transcript = None
        for lang in languages:
            try:
                transcript = transcript_list.find_transcript([lang])
                break
            except Exception:
                continue

        if not transcript:
            try:
                transcript = transcript_list.find_generated_transcript(languages)
            except Exception:
                # 嘗試抓取任意一個可用字幕
                keys = list(transcript_list._manually_created_transcripts.keys()) + list(transcript_list._generated_transcripts.keys())
                if keys:
                    transcript = transcript_list.find_transcript([keys[0]])

        if not transcript:
            raise ValueError("找不到該影片的任何字幕")

        parts = transcript.fetch()
        full_text = [part.text.strip() for part in parts if part.text.strip()]
        if not full_text:
            raise ValueError("字幕內容為空")

        return " ".join(full_text)

    def _transcribe_youtube_audio(self, url: str, video_id: str) -> str:
        """下載 YouTube 影片音軌至暫存檔，並以本地 Whisper 模型轉譯為文字"""
        with tempfile.TemporaryDirectory() as tmp_dir:
            audio_path_template = os.path.join(tmp_dir, f"{video_id}.%(ext)s")
            ydl_opts = {
                "format": "bestaudio/best",
                "outtmpl": audio_path_template,
                "quiet": True,
                "no_warnings": True,
                "noprogress": True,
            }
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.download([url])

            downloaded_files = [
                os.path.join(tmp_dir, f) for f in os.listdir(tmp_dir)
                if f.startswith(video_id)
            ]
            if not downloaded_files:
                raise DocumentParserError("音軌下載失敗，未產生任何音訊檔案")

            return _transcribe_with_whisper(downloaded_files[0], model_size="tiny")

    def _extract_youtube_id(self, url: str) -> Optional[str]:
        """從網址中擷取 YouTube Video ID"""
        patterns = [
            r'v=([^&#]+)',                  # youtube.com/watch?v=xxx
            r'youtu\.be/([^&#\?]+)',        # youtu.be/xxx
            r'embed/([^&#\?]+)',            # youtube.com/embed/xxx
            r'shorts/([^&#\?]+)'            # youtube.com/shorts/xxx
        ]
        for pattern in patterns:
            match = re.search(pattern, url)
            if match:
                return match.group(1)
        return None

    def _parse_webpage(self, url: str) -> str:
        """抓取一般網頁並轉為 Markdown"""
        if trafilatura is None:
            raise ImportError("未安裝 trafilatura 套件，請執行 pip install trafilatura")

        try:
            downloaded = trafilatura.fetch_url(url)
            if downloaded is None:
                raise ValueError(f"無法下載網頁內容: {url}")

            result = trafilatura.extract(
                downloaded,
                output_format='markdown',
                include_links=True,
                include_images=False,
                include_tables=True
            )
            
            if not result:
                result = trafilatura.extract(downloaded, output_format='txt')

            if not result:
                raise ValueError("無法提取該網頁的有效內容")

            return f"--- [網頁連結: {url}] ---\n\n" + result

        except Exception as e:
            raise DocumentParserError(f"網頁內容抓取失敗: {str(e)}")

