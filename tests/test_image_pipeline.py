import io

import pytest
from PIL import Image, ImageDraw

from parser.core import DocumentParser
from parser.image_pipeline import ImagePipeline, ImagePipelineConfig, ImageProcessResult


def _make_diagram_image(size=(300, 200)) -> Image.Image:
    """產生一張帶有方框與連線的合成圖片，模擬流程圖／架構圖的圖形特徵。"""
    img = Image.new("RGB", size, "white")
    draw = ImageDraw.Draw(img)
    draw.rectangle([20, 20, 120, 80], outline="black", width=3)
    draw.rectangle([180, 20, 280, 80], outline="black", width=3)
    draw.line([120, 50, 180, 50], fill="black", width=3)
    return img


def _make_plain_photo(size=(300, 200)) -> Image.Image:
    """產生一張無明顯結構特徵的純色圖片，模擬普通照片/裝飾圖。"""
    return Image.new("RGB", size, (120, 140, 160))


# --- 空間去重 (IoU 覆蓋率) ---------------------------------------------------

def test_spatial_dedup_skips_when_covered_by_native_text():
    pipeline = ImagePipeline()
    image_bbox = (0, 0, 100, 100)
    text_bboxes = [(0, 0, 100, 80)]  # 覆蓋 80% 面積，超過預設 70% 門檻
    assert pipeline.is_covered_by_native_text(image_bbox, text_bboxes) is True


def test_spatial_dedup_keeps_when_not_covered():
    pipeline = ImagePipeline()
    image_bbox = (0, 0, 100, 100)
    text_bboxes = [(0, 0, 100, 30)]  # 僅覆蓋 30% 面積
    assert pipeline.is_covered_by_native_text(image_bbox, text_bboxes) is False


# --- 雙序列圖號編號 ------------------------------------------------------------

def test_figure_numbering_prefers_caption_over_auto():
    pipeline = ImagePipeline()
    fid, id_type, source = pipeline.assign_figure_id("圖1 系統架構圖", "", "")
    assert fid == "圖1"
    assert id_type == "original"
    assert source == "caption"


def test_figure_numbering_falls_back_to_auto_sequence():
    pipeline = ImagePipeline()
    fid1, id_type1, _ = pipeline.assign_figure_id("", "", "")
    fid2, id_type2, _ = pipeline.assign_figure_id("", "", "")
    assert id_type1 == id_type2 == "auto"
    assert fid1 == "自補圖-1"
    assert fid2 == "自補圖-2"
    assert fid1 != fid2


def test_figure_numbering_original_and_auto_sequences_never_collide():
    pipeline = ImagePipeline()
    original_fid, _, _ = pipeline.assign_figure_id("圖1 標題", "", "")
    auto_fid, id_type, _ = pipeline.assign_figure_id("", "", "")
    assert original_fid == "圖1"
    assert auto_fid == "自補圖-1"
    assert id_type == "auto"


def test_reset_numbering_clears_state_between_documents():
    pipeline = ImagePipeline()
    pipeline.assign_figure_id("", "", "")
    pipeline.reset_numbering()
    fid, _, _ = pipeline.assign_figure_id("", "", "")
    assert fid == "自補圖-1"


# --- 強制旁路規則 ---------------------------------------------------------------

def test_forced_bypass_detects_explicit_reference_in_text():
    pipeline = ImagePipeline()
    assert pipeline.check_forced_bypass("pdf_native", "如下圖所示，系統分為三層") is True


def test_forced_bypass_respects_doc_type_force_list():
    config = ImagePipelineConfig(force_visual_parse_doc_types=["pptx"])
    pipeline = ImagePipeline(config)
    assert pipeline.check_forced_bypass("pptx", "") is True
    assert pipeline.check_forced_bypass("docx", "") is False


def test_forced_bypass_false_when_nothing_matches():
    pipeline = ImagePipeline()
    assert pipeline.check_forced_bypass("docx", "普通段落文字，沒有提到任何圖片") is False


# --- 三維度評分 ------------------------------------------------------------------

def test_score_is_within_0_100_bounds():
    pipeline = ImagePipeline()
    score = pipeline.compute_score(_make_diagram_image(), "pdf_native", "", 0.0)
    assert 0.0 <= score <= 100.0


def test_pptx_base_score_higher_than_docx_given_same_image():
    pipeline = ImagePipeline()
    img = _make_plain_photo()
    pptx_score = pipeline.compute_score(img, "pptx", "識別到的文字", 90.0)
    docx_score = pipeline.compute_score(img, "docx", "識別到的文字", 90.0)
    assert pptx_score > docx_score


def test_low_ocr_confidence_increases_score():
    pipeline = ImagePipeline()
    img = _make_plain_photo()
    low_conf_score = pipeline.compute_score(img, "pdf_native", "亂碼文字", 10.0)
    high_conf_score = pipeline.compute_score(img, "pdf_native", "清楚的文字", 95.0)
    assert low_conf_score > high_conf_score


# --- 圖片前置過濾（裝飾性小圖） -------------------------------------------------

def test_tiny_decorative_image_is_filtered_out():
    pipeline = ImagePipeline()
    tiny_icon = Image.new("RGB", (10, 10), "white")
    result = pipeline.process_image(tiny_icon, doc_type="pptx")
    assert result is None


def test_normal_sized_image_produces_result():
    pipeline = ImagePipeline()
    result = pipeline.process_image(_make_diagram_image(), doc_type="pptx")
    assert result is not None
    assert result.figure_id  # 應被賦予圖號（原生或自補）


# --- 圖理解模型預設關閉、優雅降級 -------------------------------------------------

def test_image_understanding_disabled_by_default_keeps_ocr_only():
    config = ImagePipelineConfig()
    assert config.enable_image_understanding is False
    pipeline = ImagePipeline(config)
    result = pipeline.process_image(_make_diagram_image(), doc_type="pptx", force_image_understanding=False)
    assert result is not None
    assert result.used_image_understanding is False
    assert result.understanding_text is None


def test_vision_call_gracefully_degrades_when_ollama_unreachable():
    """即使強制觸發圖理解，Ollama 服務未啟動時也不應拋出例外，應自動降級保留 OCR。"""
    config = ImagePipelineConfig(
        enable_image_understanding=True,
        ollama_base_url="http://localhost:1",  # 刻意指向不存在的服務
        ollama_timeout_seconds=2,
    )
    pipeline = ImagePipeline(config)
    result = pipeline.process_image(
        _make_diagram_image(), doc_type="pptx", force_image_understanding=True
    )
    assert result is not None
    assert result.understanding_text is None
    assert result.used_image_understanding is False


class _FakeResponse:
    def __init__(self, json_data, status_code=200):
        self._json_data = json_data
        self.status_code = status_code

    def raise_for_status(self):
        if self.status_code >= 400:
            raise Exception(f"HTTP {self.status_code}")

    def json(self):
        return self._json_data


def test_vision_availability_checked_only_once_per_document(monkeypatch):
    """迴歸測試：同一份文件處理多張圖片時，Ollama 可用性只應實際檢查一次
    （打 /api/tags），不應每張符合條件的圖片都各自重新嘗試連線。"""
    pytest.importorskip("requests")
    import parser.image_pipeline as ip_module

    get_calls = []

    def fake_get(url, timeout=None):
        get_calls.append(url)
        return _FakeResponse({"models": [{"name": "qwen2.5vl:7b"}]})

    monkeypatch.setattr(ip_module.requests, "get", fake_get)
    monkeypatch.setattr(
        ip_module.requests, "post",
        lambda *a, **k: _FakeResponse({"response": "一張示意圖"}),
    )

    config = ImagePipelineConfig(enable_image_understanding=True)
    pipeline = ImagePipeline(config)

    for _ in range(3):
        result = pipeline.process_image(
            _make_diagram_image(), doc_type="pptx", force_image_understanding=True
        )
        assert result.used_image_understanding is True

    assert len(get_calls) == 1  # 只打了一次 /api/tags，不是三次


def test_vision_availability_check_detects_model_not_installed(monkeypatch):
    """迴歸測試：Ollama 服務可連線，但指定的模型不在已安裝清單中時，
    應判定為不可用並優雅降級，而不是直接嘗試呼叫 /api/generate 才發現失敗。"""
    pytest.importorskip("requests")
    import parser.image_pipeline as ip_module

    post_calls = []
    monkeypatch.setattr(
        ip_module.requests, "get",
        lambda url, timeout=None: _FakeResponse({"models": [{"name": "llava:latest"}]}),
    )
    monkeypatch.setattr(
        ip_module.requests, "post",
        lambda *a, **k: post_calls.append(1) or _FakeResponse({"response": "不應被呼叫"}),
    )

    config = ImagePipelineConfig(
        enable_image_understanding=True, ollama_vision_model="qwen2.5vl:7b",
    )
    pipeline = ImagePipeline(config)
    result = pipeline.process_image(
        _make_diagram_image(), doc_type="pptx", force_image_understanding=True
    )

    assert result.used_image_understanding is False
    assert result.understanding_text is None
    assert post_calls == []  # 健康檢查判定不可用後，不應再嘗試呼叫 /api/generate


def test_vision_availability_check_matches_model_family_ignoring_tag_variant(monkeypatch):
    """模型清單裡的 tag 可能跟設定值不完全一致（如 qwen2.5vl:latest vs qwen2.5vl:7b），
    只要模型家族名稱前綴相符，仍應判定為可用。"""
    pytest.importorskip("requests")
    import parser.image_pipeline as ip_module

    monkeypatch.setattr(
        ip_module.requests, "get",
        lambda url, timeout=None: _FakeResponse({"models": [{"name": "qwen2.5vl:latest"}]}),
    )
    monkeypatch.setattr(
        ip_module.requests, "post",
        lambda *a, **k: _FakeResponse({"response": "描述文字"}),
    )

    config = ImagePipelineConfig(
        enable_image_understanding=True, ollama_vision_model="qwen2.5vl:7b",
    )
    pipeline = ImagePipeline(config)
    result = pipeline.process_image(
        _make_diagram_image(), doc_type="pptx", force_image_understanding=True
    )

    assert result.used_image_understanding is True


def test_reset_numbering_allows_rechecking_vision_availability(monkeypatch):
    """迴歸測試：reset_numbering()（每份新文件開始時呼叫）應重置可用性快取，
    讓下一份文件有機會重新確認 Ollama 狀態，不會被上一份文件的檢查結果卡住。"""
    pytest.importorskip("requests")
    import parser.image_pipeline as ip_module

    get_calls = []

    def fake_get(url, timeout=None):
        get_calls.append(url)
        return _FakeResponse({"models": [{"name": "qwen2.5vl:7b"}]})

    monkeypatch.setattr(ip_module.requests, "get", fake_get)
    monkeypatch.setattr(
        ip_module.requests, "post",
        lambda *a, **k: _FakeResponse({"response": "描述文字"}),
    )

    config = ImagePipelineConfig(enable_image_understanding=True)
    pipeline = ImagePipeline(config)

    pipeline.process_image(_make_diagram_image(), doc_type="pptx", force_image_understanding=True)
    pipeline.process_image(_make_diagram_image(), doc_type="pptx", force_image_understanding=True)
    assert len(get_calls) == 1

    pipeline.reset_numbering()
    pipeline.process_image(_make_diagram_image(), doc_type="pptx", force_image_understanding=True)
    assert len(get_calls) == 2  # 重置後，新文件重新檢查了一次


# --- 輸出格式化 ------------------------------------------------------------------

def test_to_text_block_includes_figure_id_and_content():
    result = ImageProcessResult(
        figure_id="圖1", id_type="original", match_source="caption",
        caption="系統架構圖", ocr_text="輸入 輸出", ocr_confidence=90.0, score=20.0,
    )
    block = result.to_text_block()
    assert "[圖片:圖1]" in block
    assert "系統架構圖" in block
    assert "輸入 輸出" in block


def test_to_text_block_prefers_understanding_over_ocr():
    result = ImageProcessResult(
        figure_id="自補圖-1", id_type="auto", match_source="none",
        ocr_text="模糊文字", score=80.0,
        used_image_understanding=True, understanding_text="流程圖：A 指向 B",
    )
    block = result.to_text_block()
    assert "流程圖：A 指向 B" in block
    assert "模糊文字" not in block


# --- 端對端：PPTX 內嵌圖片會被抽取並回填進最終輸出文字 ---------------------------

def test_pptx_end_to_end_extracts_embedded_image_content(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面

    img_buf = io.BytesIO()
    _make_diagram_image().save(img_buf, format="PNG")
    img_buf.seek(0)
    slide.shapes.add_picture(img_buf, left=0, top=0)

    pptx_path = tmp_path / "sample.pptx"
    prs.save(pptx_path)

    parser = DocumentParser()
    text = parser.parse_file(str(pptx_path))

    assert "[圖片:" in text


def test_docx_end_to_end_extracts_embedded_image_content(tmp_path):
    docx_module = pytest.importorskip("docx")
    from docx import Document

    doc = Document()
    doc.add_paragraph("這是段落文字，下方附有架構圖。")
    img_buf = io.BytesIO()
    _make_diagram_image().save(img_buf, format="PNG")
    img_buf.seek(0)
    doc.add_picture(img_buf)

    docx_path = tmp_path / "sample.docx"
    doc.save(docx_path)

    parser = DocumentParser()
    text = parser.parse_file(str(docx_path))

    assert "[圖片:" in text
    assert "這是段落文字" in text


def test_pdf_end_to_end_extracts_embedded_image_content(tmp_path):
    """驗證 PDF 內嵌 JPEG 圖片能被 pdfminer LTImage 解碼並經圖片管線回填進輸出文字。"""
    reportlab = pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader

    img_buf = io.BytesIO()
    _make_diagram_image().save(img_buf, format="JPEG")
    img_buf.seek(0)

    pdf_path = tmp_path / "sample.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=(400, 500))
    # 使用英文內容：reportlab 內建 Helvetica 字型不含 CJK 字形，避免測試混入字型渲染問題。
    # 需要足夠長度的原生文字，避免觸發低品質判定而改走 OCR 備援軌道（該軌道需要 Poppler）。
    c.drawString(50, 460, "Body paragraph describing the system, diagram is shown below.")
    c.drawString(50, 440, "This filler text ensures native text extraction passes the")
    c.drawString(50, 420, "length and readable-character-ratio quality thresholds.")
    c.drawImage(ImageReader(img_buf), 50, 100, width=300, height=200)
    c.save()

    parser = DocumentParser()
    text = parser.parse_file(str(pdf_path))

    assert "Body paragraph" in text
    assert "[圖片:" in text


def test_pdf_image_numbering_follows_visual_order_not_draw_order(tmp_path):
    """迴歸測試：自補圖號的指派順序應依頁面視覺由上到下排序，
    不應受 pdfminer 內部圖形元素發現順序（近似 PDF content stream 繪製順序）影響。"""
    reportlab = pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader

    # 以不同尺寸區分兩張圖片，供測試辨識呼叫順序對應到哪一張
    top_img_buf = io.BytesIO()
    Image.new("RGB", (300, 80), (200, 200, 200)).save(top_img_buf, format="JPEG")
    top_img_buf.seek(0)

    bottom_img_buf = io.BytesIO()
    Image.new("RGB", (300, 150), (100, 100, 100)).save(bottom_img_buf, format="JPEG")
    bottom_img_buf.seek(0)

    pdf_path = tmp_path / "order_test.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=(400, 500))
    c.drawString(50, 470, "Body text ensures native extraction passes quality thresholds for this test document reliably.")
    c.drawString(50, 450, "Additional filler sentence to keep the readable character ratio comfortably above the cutoff.")
    # 刻意先畫「視覺上位於下方」的圖片，模擬 pdfminer 內部發現順序與視覺閱讀順序不同的情境
    c.drawImage(ImageReader(bottom_img_buf), 50, 80, width=300, height=150)   # y: 80~230（下方）
    c.drawImage(ImageReader(top_img_buf), 50, 260, width=300, height=80)     # y: 260~340（上方）
    c.save()

    parser = DocumentParser()
    call_order_sizes = []
    original_process_image = parser.image_pipeline.process_image

    def recording_process_image(pil_image, **kwargs):
        call_order_sizes.append(pil_image.size)
        return original_process_image(pil_image, **kwargs)

    parser.image_pipeline.process_image = recording_process_image

    text = parser.parse_file(str(pdf_path))

    assert len(call_order_sizes) == 2
    # 視覺上方的圖（300x80）應先被指派圖號，即使它在 content stream 中是後畫的
    assert call_order_sizes[0] == (300, 80)
    assert call_order_sizes[1] == (300, 150)

    assert "自補圖-1" in text and "自補圖-2" in text
    assert text.index("自補圖-1") < text.index("自補圖-2")


# --- 純文字文件預檢：跳過整套圖片管線 -----------------------------------------

def _make_text_only_pdf(tmp_path, filename="text_only.pdf"):
    reportlab = pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas

    pdf_path = tmp_path / filename
    c = canvas.Canvas(str(pdf_path), pagesize=(400, 500))
    c.drawString(50, 460, "This document contains only native text and no embedded images at all.")
    c.drawString(50, 440, "It exists purely to validate that the cheap pre-check skips the image pipeline.")
    c.drawString(50, 420, "A third filler line keeps the readable-character-ratio comfortably above threshold.")
    c.save()
    return pdf_path


def test_pdf_has_embedded_images_false_for_text_only_pdf(tmp_path):
    pdf_path = _make_text_only_pdf(tmp_path)
    parser = DocumentParser()
    assert parser._pdf_has_embedded_images(pdf_path) is False


def test_pdf_has_embedded_images_true_when_image_present(tmp_path):
    reportlab = pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader

    img_buf = io.BytesIO()
    _make_diagram_image().save(img_buf, format="JPEG")
    img_buf.seek(0)

    pdf_path = tmp_path / "has_image.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=(400, 500))
    c.drawString(50, 460, "Body text plus an embedded diagram below it for this test document.")
    c.drawImage(ImageReader(img_buf), 50, 100, width=300, height=200)
    c.save()

    parser = DocumentParser()
    assert parser._pdf_has_embedded_images(pdf_path) is True


def test_pdf_text_only_skips_image_pipeline_entirely(tmp_path, monkeypatch):
    """迴歸測試：純文字 PDF 不應觸發 _extract_pdf_images_text（避免重複 extract_pages()）。"""
    pdf_path = _make_text_only_pdf(tmp_path)
    parser = DocumentParser()

    called = []
    original = parser._extract_pdf_images_text

    def spy(path):
        called.append(path)
        return original(path)

    monkeypatch.setattr(parser, "_extract_pdf_images_text", spy)

    text = parser.parse_file(str(pdf_path))

    assert called == []
    assert "[圖片:" not in text
    assert "This document contains only native text" in text


def test_pdf_with_image_still_invokes_image_pipeline(tmp_path, monkeypatch):
    """對照組：含圖片的 PDF 仍應正常觸發圖片管線，確認預檢沒有誤傷正常情境。"""
    reportlab = pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader

    img_buf = io.BytesIO()
    _make_diagram_image().save(img_buf, format="JPEG")
    img_buf.seek(0)

    pdf_path = tmp_path / "with_image.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=(400, 500))
    c.drawString(50, 460, "Body text plus an embedded diagram shown further below on this page.")
    c.drawString(50, 440, "This filler line ensures native text extraction passes the quality gate.")
    c.drawString(50, 420, "A third line keeps the readable-character-ratio comfortably above cutoff.")
    c.drawImage(ImageReader(img_buf), 50, 100, width=300, height=200)
    c.save()

    parser = DocumentParser()

    called = []
    original = parser._extract_pdf_images_text

    def spy(path):
        called.append(path)
        return original(path)

    monkeypatch.setattr(parser, "_extract_pdf_images_text", spy)

    text = parser.parse_file(str(pdf_path))

    assert len(called) == 1
    assert "[圖片:" in text


def test_docx_has_embedded_images_false_for_text_only_docx(tmp_path):
    docx_module = pytest.importorskip("docx")
    from docx import Document

    doc = Document()
    doc.add_paragraph("這是一份完全沒有圖片的純文字文件。")
    doc.add_paragraph("第二段落，用來確認純文字偵測邏輯正確運作。")
    docx_path = tmp_path / "text_only.docx"
    doc.save(docx_path)

    parser = DocumentParser()
    reopened = Document(docx_path)
    assert parser._docx_has_embedded_images(reopened) is False


def test_docx_text_only_skips_paragraph_image_scan(tmp_path, monkeypatch):
    """迴歸測試：純文字 DOCX 不應對任何段落呼叫 _process_docx_paragraph_images。"""
    docx_module = pytest.importorskip("docx")
    from docx import Document

    doc = Document()
    for i in range(5):
        doc.add_paragraph(f"第 {i + 1} 段純文字內容，沒有任何內嵌圖片。")
    docx_path = tmp_path / "text_only_multi.docx"
    doc.save(docx_path)

    parser = DocumentParser()
    called = []
    original = parser._process_docx_paragraph_images

    def spy(doc_obj, paragraph, full_doc_text):
        called.append(paragraph)
        return original(doc_obj, paragraph, full_doc_text)

    monkeypatch.setattr(parser, "_process_docx_paragraph_images", spy)

    text = parser.parse_file(str(docx_path))

    assert called == []
    assert "[圖片:" not in text
    assert "第 1 段純文字內容" in text


def test_docx_with_image_still_invokes_paragraph_image_scan(tmp_path, monkeypatch):
    """對照組：含圖片的 DOCX 仍應正常掃描段落圖片，確認預檢沒有誤傷正常情境。"""
    docx_module = pytest.importorskip("docx")
    from docx import Document

    doc = Document()
    doc.add_paragraph("這段文字後面有一張圖片。")
    img_buf = io.BytesIO()
    _make_diagram_image().save(img_buf, format="PNG")
    img_buf.seek(0)
    doc.add_picture(img_buf)
    docx_path = tmp_path / "with_image.docx"
    doc.save(docx_path)

    parser = DocumentParser()
    called = []
    original = parser._process_docx_paragraph_images

    def spy(doc_obj, paragraph, full_doc_text):
        called.append(paragraph)
        return original(doc_obj, paragraph, full_doc_text)

    monkeypatch.setattr(parser, "_process_docx_paragraph_images", spy)

    text = parser.parse_file(str(docx_path))

    assert len(called) >= 1
    assert "[圖片:" in text


# --- 合併另一 session 優化後的迴歸測試 -----------------------------------------

def test_pptx_has_embedded_images_true_for_real_pptx_with_picture(tmp_path):
    """迴歸測試：PPTX 圖片關聯存在於各 slide 自己的 rels，並非簡報層級的 prs.part.rels。
    這裡直接針對真實 python-pptx 物件驗證，避免用 mock 掩蓋掉這個結構性差異。"""
    pptx_module = pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img_buf = io.BytesIO()
    _make_diagram_image().save(img_buf, format="PNG")
    img_buf.seek(0)
    slide.shapes.add_picture(img_buf, left=0, top=0)

    parser = DocumentParser()
    assert parser._pptx_has_embedded_images(prs) is True


def test_pptx_has_embedded_images_false_for_text_only_pptx(tmp_path):
    pptx_module = pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 100, 100).text_frame.text = "純文字投影片，沒有任何圖片。"

    parser = DocumentParser()
    assert parser._pptx_has_embedded_images(prs) is False


def test_pptx_text_only_skips_picture_shape_processing(tmp_path, monkeypatch):
    pptx_module = pytest.importorskip("pptx")
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(0, 0, 100, 100).text_frame.text = "純文字投影片，完全沒有圖片。"
    pptx_path = tmp_path / "text_only.pptx"
    prs.save(pptx_path)

    parser = DocumentParser()
    called = []
    original = parser._process_pptx_picture_shape

    def spy(shape, nearby_text, full_slide_text):
        called.append(shape)
        return original(shape, nearby_text, full_slide_text)

    monkeypatch.setattr(parser, "_process_pptx_picture_shape", spy)

    text = parser.parse_file(str(pptx_path))

    assert called == []
    assert "[圖片:" not in text
    assert "純文字投影片" in text


def test_pptx_reading_order_follows_visual_position_not_z_order():
    """迴歸測試：_pptx_reading_order 應依視覺 (top, left) 排序，
    不受 shapes 加入順序（z-order）影響。"""

    class FakeShape:
        def __init__(self, name, top, left):
            self.name = name
            self.top = top
            self.left = left

    # 刻意以「視覺上在下方/右方」的形狀先加入清單，模擬 z-order 與視覺順序不同的情境
    bottom_right = FakeShape("bottom_right", top=500000, left=500000)
    top_left = FakeShape("top_left", top=0, left=0)
    top_right = FakeShape("top_right", top=0, left=500000)  # 與 top_left 同列（容忍度內）

    parser = DocumentParser()
    ordered = parser._pptx_reading_order([bottom_right, top_right, top_left])

    assert [s.name for s in ordered] == ["top_left", "top_right", "bottom_right"]


def test_docx_caption_found_in_next_paragraph_when_image_paragraph_is_empty(tmp_path):
    """迴歸測試：Word「插入標題」慣例——圖片自成一段（無文字），標題在下一段——
    應能正確抓取原文圖號，而非落入自補序列。"""
    docx_module = pytest.importorskip("docx")
    from docx import Document

    doc = Document()
    doc.add_paragraph("前言段落，說明後續會有一張架構圖。")
    img_buf = io.BytesIO()
    _make_diagram_image().save(img_buf, format="PNG")
    img_buf.seek(0)
    doc.add_picture(img_buf)  # 圖片自成一段，本身沒有文字
    doc.add_paragraph("圖1 系統架構圖")  # 緊接的下一段才是標題

    docx_path = tmp_path / "caption_next_paragraph.docx"
    doc.save(docx_path)

    parser = DocumentParser()
    text = parser.parse_file(str(docx_path))

    assert "[圖片:圖1]" in text
    assert "自補圖" not in text


def test_pdf_pages_layout_shared_between_track_two_and_image_pipeline(tmp_path, monkeypatch):
    """迴歸測試：當文字提取走到軌道二、且文件確實含有內嵌圖片時，
    pdfminer 的版面樹狀結構應只載入一次（_load_pdf_layout 只被呼叫一次），
    供軌道二文字重建與圖片管線共用，而非各自呼叫一次 extract_pages()。"""
    reportlab = pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader

    img_buf = io.BytesIO()
    _make_diagram_image().save(img_buf, format="JPEG")
    img_buf.seek(0)

    pdf_path = tmp_path / "cache_test.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=(400, 500))
    c.drawString(50, 460, "Body text long enough to satisfy the quality gate once reconstructed by pdfminer.")
    c.drawString(50, 440, "A second filler line keeps the readable-character-ratio comfortably above cutoff.")
    c.drawImage(ImageReader(img_buf), 50, 100, width=300, height=200)
    c.save()

    parser = DocumentParser()

    load_calls = []
    original_load = parser._load_pdf_layout

    def spy_load(path):
        load_calls.append(path)
        return original_load(path)

    monkeypatch.setattr(parser, "_load_pdf_layout", spy_load)

    quality_calls = {"n": 0}
    original_quality = parser._is_low_quality_text

    def force_track_two_only(text):
        quality_calls["n"] += 1
        # 第一次呼叫是軌道一（pypdf）輸出，強制判定為低品質以逼迫走軌道二；
        # 之後的呼叫交由原始邏輯判斷，讓軌道二重建出的正常文字通過品質檢查、
        # 不再繼續掉到軌道三（否則圖片管線會因 used_ocr_fallback 而被跳過，
        # 就測不到「軌道二與圖片管線共用同一份版面樹」這件事了）。
        if quality_calls["n"] == 1:
            return True
        return original_quality(text)

    monkeypatch.setattr(parser, "_is_low_quality_text", force_track_two_only)

    text = parser.parse_file(str(pdf_path))

    assert len(load_calls) == 1
    assert "[圖片:" in text


def test_pdf_reader_none_still_falls_back_to_pdfminer_track(tmp_path, monkeypatch):
    """迴歸測試：即使 pypdf 開檔失敗（或未安裝），仍應優雅降級走軌道二，
    而非直接回傳空字串，中斷整個 PDF 解析。"""
    reportlab = pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas

    pdf_path = tmp_path / "reader_none_fallback.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=(400, 500))
    c.drawString(50, 460, "Body text that must still be recovered via pdfminer even when pypdf is unavailable.")
    c.drawString(50, 440, "A second filler line keeps the readable-character-ratio comfortably above cutoff.")
    c.save()

    import parser.core as core_module
    monkeypatch.setattr(core_module, "PdfReader", None)

    parser = DocumentParser()
    text = parser.parse_file(str(pdf_path))

    assert "Body text that must still be recovered" in text


# --- check_forced_bypass 三層邏輯（精確圖號比對 + 無編號安全網） -----------------

def test_forced_bypass_original_id_triggers_on_specific_reference():
    pipeline = ImagePipeline()
    assert pipeline.check_forced_bypass(
        "pdf_native", "", figure_id="圖3", id_type="original",
        full_doc_text="前文省略。如圖3所示，系統分為三層架構。",
    ) is True


def test_forced_bypass_original_id_does_not_cross_trigger_other_figure():
    """迴歸測試：正文只精確引用圖2，不應連帶誤觸發圖3（避免同頁多圖過度觸發）。
    注意：counter-example 文字刻意不得包含「圖3」這個子字串本身，
    否則會因為比對 pattern 的前後綴皆為可選而誤判為命中。"""
    pipeline = ImagePipeline()
    assert pipeline.check_forced_bypass(
        "pdf_native", "", figure_id="圖3", id_type="original",
        full_doc_text="前文省略。如圖2所示，系統分為三層架構，此段完全沒有提到其他圖表。",
    ) is False


def test_forced_bypass_auto_id_falls_back_to_generic_pattern():
    """迴歸測試：無原文圖號可比對（自補序列）時，仍應保留通用旁路關鍵字的安全網，
    避免「單一無編號圖片＋泛用引用語句」的情境完全漏判。"""
    pipeline = ImagePipeline()
    assert pipeline.check_forced_bypass(
        "pdf_native", "", figure_id="自補圖-1", id_type="auto",
        full_doc_text="如下圖所示，系統分為三層架構，圖中未標號。",
    ) is True


def test_forced_bypass_auto_id_without_generic_phrase_stays_false():
    pipeline = ImagePipeline()
    assert pipeline.check_forced_bypass(
        "pdf_native", "", figure_id="自補圖-1", id_type="auto",
        full_doc_text="這是一段完全沒有提到任何圖片的普通文字。",
    ) is False


# --- 裝飾小圖過濾：改用顯示尺寸而非原始像素尺寸 -------------------------------

def test_unit_conversion_helpers():
    from parser.image_pipeline import emu_to_px, points_to_px

    assert emu_to_px(914400) == pytest.approx(96.0)   # 1 inch = 914400 EMU = 96px (96 DPI)
    assert points_to_px(72) == pytest.approx(96.0)     # 1 inch = 72 point = 96px (96 DPI)


def test_process_image_uses_display_size_over_raw_pixel_size_high_res_small_display():
    """迴歸測試：高解析度原圖但顯示尺寸很小（例如縮小成裝飾用小圖示），
    應依顯示尺寸判定為裝飾性小圖並濾除，而非因原始像素夠大而誤放行。"""
    pipeline = ImagePipeline()
    high_res_image = Image.new("RGB", (3000, 2000), "white")  # 原始像素遠大於 40px

    result = pipeline.process_image(
        high_res_image, doc_type="pptx", display_size_px=(28.0, 28.0),  # 顯示尺寸 <40px
    )
    assert result is None


def test_process_image_uses_display_size_over_raw_pixel_size_low_res_large_display():
    """迴歸測試：原始像素很小但顯示尺寸放大到有意義的大小，
    不應僅因原始解析度低於 40px 就被誤判為裝飾性小圖而濾除。"""
    pipeline = ImagePipeline()
    low_res_image = _make_plain_photo(size=(35, 35))  # 原始像素 <40px

    result = pipeline.process_image(
        low_res_image, doc_type="pptx", display_size_px=(384.0, 288.0),  # 顯示尺寸夠大
    )
    assert result is not None


def test_process_image_falls_back_to_raw_pixel_size_when_display_size_unavailable():
    """未提供 display_size_px（呼叫端無法取得顯示尺寸）時，應退回使用原始像素尺寸判斷，
    維持向後相容，不因無法取得顯示尺寸而讓所有圖片都被濾除或都放行。"""
    pipeline = ImagePipeline()
    tiny_image = Image.new("RGB", (10, 10), "white")

    assert pipeline.process_image(tiny_image, doc_type="pptx") is None
    assert pipeline.process_image(_make_diagram_image(), doc_type="pptx") is not None


def test_pdf_extraction_computes_display_size_from_bbox_in_points(tmp_path):
    """端對端驗證：PDF 圖片的 bbox（point）會被正確換算成顯示像素，一張顯示尺寸
    很小（雖然原始解碼像素不小）的圖應該被濾除，不出現在最終輸出中。"""
    reportlab = pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader

    # 原始像素 300x200（遠大於 40px），但在頁面上只用 20x20 point 顯示
    # 20pt 換算成 96 DPI px ≈ 26.7px，低於預設門檻 40px
    img_buf = io.BytesIO()
    _make_diagram_image().save(img_buf, format="JPEG")
    img_buf.seek(0)

    pdf_path = tmp_path / "tiny_display.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=(400, 500))
    c.drawString(50, 460, "Body text long enough to satisfy the native text quality gate reliably.")
    c.drawString(50, 440, "A second filler line keeps the readable-character-ratio above the cutoff.")
    c.drawImage(ImageReader(img_buf), 50, 100, width=20, height=20)
    c.save()

    parser = DocumentParser()
    text = parser.parse_file(str(pdf_path))

    assert "[圖片:" not in text


def test_pptx_extraction_computes_display_size_from_shape_emu(tmp_path):
    """端對端驗證：PPTX 圖片 shape 的 width/height（EMU）會被正確換算成顯示像素，
    一張原始像素很大但顯示很小的圖應該被濾除。"""
    pptx_module = pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img_buf = io.BytesIO()
    Image.new("RGB", (3000, 2000), "white").save(img_buf, format="PNG")
    img_buf.seek(0)
    # 顯式指定顯示尺寸為 20x20 EMU px 換算下遠小於 40px（約 9525*20 EMU ≈ 20px）
    slide.shapes.add_picture(img_buf, left=0, top=0, width=Emu(190500), height=Emu(190500))

    pptx_path = tmp_path / "tiny_display.pptx"
    prs.save(pptx_path)

    parser = DocumentParser()
    text = parser.parse_file(str(pptx_path))

    assert "[圖片:" not in text


def test_docx_extraction_computes_display_size_from_wp_extent(tmp_path):
    """端對端驗證：DOCX 圖片的 wp:extent（EMU，顯示尺寸）會被正確換算成顯示像素，
    一張原始像素很大但顯示很小（例如縮成裝飾用小圖示）的圖應該被濾除。"""
    docx_module = pytest.importorskip("docx")
    from docx import Document
    from docx.shared import Emu

    doc = Document()
    p = doc.add_paragraph()
    run = p.add_run()
    img_buf = io.BytesIO()
    Image.new("RGB", (3000, 2000), "white").save(img_buf, format="PNG")
    img_buf.seek(0)
    # 顯示尺寸縮小到約 20px（190500 EMU / 9525 = 20px）
    run.add_picture(img_buf, width=Emu(190500), height=Emu(190500))

    docx_path = tmp_path / "tiny_display.docx"
    doc.save(docx_path)

    parser = DocumentParser()
    text = parser.parse_file(str(docx_path))

    assert "[圖片:" not in text


def test_docx_extraction_keeps_low_res_image_with_large_display_size(tmp_path):
    """對照組：原始像素小但顯示尺寸放大的 DOCX 圖片不應被誤判為裝飾性小圖。"""
    docx_module = pytest.importorskip("docx")
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    p = doc.add_paragraph()
    run = p.add_run()
    img_buf = io.BytesIO()
    _make_plain_photo(size=(35, 35)).save(img_buf, format="PNG")  # 原始像素 <40px
    img_buf.seek(0)
    run.add_picture(img_buf, width=Inches(3), height=Inches(3))  # 顯示尺寸夠大

    docx_path = tmp_path / "large_display.docx"
    doc.save(docx_path)

    parser = DocumentParser()
    text = parser.parse_file(str(docx_path))

    assert "[圖片:" in text
